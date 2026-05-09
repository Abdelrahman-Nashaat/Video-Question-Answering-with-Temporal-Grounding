"""Build the project presentation deck (presentation.pptx).

Run this script from the repo root:
    python deliverables/build_presentation.py

Produces deliverables/presentation.pptx (16:9, 11 slides).
Pure-python, depends only on python-pptx.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ACCENT = RGBColor(0x1E, 0x40, 0xAF)
TEXT_DARK = RGBColor(0x10, 0x18, 0x27)
TEXT_MUTED = RGBColor(0x4B, 0x55, 0x63)
BG_LIGHT = RGBColor(0xFF, 0xFF, 0xFF)

TITLE_FONT_SIZE = Pt(36)
BULLET_FONT_SIZE = Pt(22)
CODE_FONT_SIZE = Pt(14)
FOOTER_FONT_SIZE = Pt(11)

DEFAULT_FONT = "Calibri"
MONO_FONT = "Courier New"

ARCHITECTURE_DIAGRAM = r"""                 +----------------------------------------------+
                 |  Inputs: video (file or URL) + question      |
                 +------------------------+---------------------+
                                          |
              +---------------------------+---------------------------+
              v                                                       v
   +--------------------+                                  +----------------------+
   | Frame extraction   |                                  | Audio extraction     |
   |  - 2 fps sampling  |                                  |  - faster-whisper    |
   |  - histogram dedup |                                  |    large-v3 (auto    |
   +---------+----------+                                  |    language detect)  |
             v                                             +-----------+----------+
   +--------------------+                                              v
   | SigLIP image embed |                                  +----------------------+
   | (large-384 / base) |                                  | Sliding-window chunks|
   +---------+----------+                                  |  + bge-m3 embeddings |
             |                                             +-----------+----------+
             +-----------------------+----------------------------------+
                                     v
            +----------------------------------------------------+
            | Per-timestamp raw scores -> z-softmax normalize    |
            | -> alpha*v_norm + beta*a_norm -> top-K moments     |
            +-----------------------+----------------------------+
                                    v
            +----------------------------------------------------+
            | Llama 3.3 70B (Groq) -> grounded answer +          |
            | cited [mm:ss] timestamps                           |
            +-----------------------+----------------------------+
                                    v
            +----------------------------------------------------+
            | XAI: timeline . Grad-CAM . frame grid . modality   |
            +-----------------------+----------------------------+
                                    v
                              +-------------+
                              |  Gradio UI  |
                              |  share=True |
                              +-------------+"""


def add_top_bar(slide, prs):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.18)
    )
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.shadow.inherit = False


def add_footer(slide, prs, page_number, total):
    box = slide.shapes.add_textbox(
        Inches(0.4),
        prs.slide_height - Inches(0.45),
        prs.slide_width - Inches(0.8),
        Inches(0.35),
    )
    tf = box.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{page_number} / {total}"
    run.font.name = DEFAULT_FONT
    run.font.size = FOOTER_FONT_SIZE
    run.font.color.rgb = TEXT_MUTED


def add_textbox(slide, left, top, width, height, text, *, font_size,
                bold=False, color=TEXT_DARK, font=DEFAULT_FONT,
                align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, bullets, *,
                font_size=BULLET_FONT_SIZE, color=TEXT_DARK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = "•  " + item
        run.font.name = DEFAULT_FONT
        run.font.size = font_size
        run.font.color.rgb = color
    return box


def add_slide_title(slide, prs, text):
    add_textbox(
        slide,
        Inches(0.6),
        Inches(0.45),
        prs.slide_width - Inches(1.2),
        Inches(0.9),
        text,
        font_size=TITLE_FONT_SIZE,
        bold=True,
        color=TEXT_DARK,
    )


def blank_slide(prs):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_LIGHT
    return slide


def slide_title(prs, total):
    slide = blank_slide(prs)
    add_top_bar(slide, prs)
    accent_block = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.0),
        Inches(2.2),
        prs.slide_width,
        Inches(2.6),
    )
    accent_block.line.fill.background()
    accent_block.fill.solid()
    accent_block.fill.fore_color.rgb = ACCENT
    accent_block.shadow.inherit = False

    add_textbox(
        slide,
        Inches(0.6),
        Inches(2.45),
        prs.slide_width - Inches(1.2),
        Inches(1.2),
        "Video Question Answering with Temporal Grounding",
        font_size=Pt(40),
        bold=True,
        color=BG_LIGHT,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        Inches(0.6),
        Inches(3.65),
        prs.slide_width - Inches(1.2),
        Inches(0.6),
        "A Multi-modal System with Explainable AI",
        font_size=Pt(22),
        color=BG_LIGHT,
        align=PP_ALIGN.CENTER,
    )

    add_textbox(
        slide,
        Inches(0.6),
        Inches(5.3),
        prs.slide_width - Inches(1.2),
        Inches(0.5),
        "Course:  Computer Vision",
        font_size=Pt(18),
        bold=True,
        color=TEXT_DARK,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        Inches(0.6),
        Inches(5.85),
        prs.slide_width - Inches(1.2),
        Inches(0.5),
        "Author:  Abdelrahman Nashaan",
        font_size=Pt(18),
        color=TEXT_DARK,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        Inches(0.6),
        Inches(6.35),
        prs.slide_width - Inches(1.2),
        Inches(0.5),
        "Tanta University — Faculty of Artificial Intelligence",
        font_size=Pt(16),
        color=TEXT_MUTED,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, prs, 1, total)


def slide_problem(prs, total):
    slide = blank_slide(prs)
    add_top_bar(slide, prs)
    add_slide_title(slide, prs, "Problem Statement")
    add_bullets(
        slide,
        Inches(0.7),
        Inches(1.7),
        prs.slide_width - Inches(1.4),
        Inches(5.0),
        [
            "Videos are dense, time-consuming sources of information.",
            "Users often need a specific answer, not the full video.",
            "Existing tools either don't localize the answer in time or don't explain how they found it.",
            "Goal — answer a free-form question about a video AND show where and why.",
        ],
    )
    add_footer(slide, prs, 2, total)


def slide_goals(prs, total):
    slide = blank_slide(prs)
    add_top_bar(slide, prs)
    add_slide_title(slide, prs, "Project Goals")
    add_bullets(
        slide,
        Inches(0.7),
        Inches(1.7),
        prs.slide_width - Inches(1.4),
        Inches(5.0),
        [
            "Multi-modal: fuse visual frames + audio transcript + question text.",
            "Bilingual: support Arabic and English questions and audio.",
            "Explainable: surface model reasoning, not just outputs.",
            "Reproducible: single notebook, runs on free Colab GPU.",
        ],
    )
    add_footer(slide, prs, 3, total)


def slide_architecture(prs, total):
    slide = blank_slide(prs)
    add_top_bar(slide, prs)
    add_slide_title(slide, prs, "System Architecture")

    box = slide.shapes.add_textbox(
        Inches(0.35),
        Inches(1.45),
        prs.slide_width - Inches(0.7),
        Inches(5.2),
    )
    tf = box.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = ARCHITECTURE_DIAGRAM
    run.font.name = MONO_FONT
    run.font.size = Pt(9)
    run.font.color.rgb = TEXT_DARK

    add_textbox(
        slide,
        Inches(0.6),
        Inches(6.75),
        prs.slide_width - Inches(1.2),
        Inches(0.5),
        "Inputs → frame & audio extraction → multi-modal embedding "
        "→ retrieval → LLM answer → 4 XAI views",
        font_size=Pt(14),
        color=TEXT_MUTED,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, prs, 4, total)


def slide_pipeline(prs, total):
    slide = blank_slide(prs)
    add_top_bar(slide, prs)
    add_slide_title(slide, prs, "Multi-modal Pipeline")
    add_bullets(
        slide,
        Inches(0.7),
        Inches(1.7),
        prs.slide_width - Inches(1.4),
        Inches(5.0),
        [
            "Visual stream: SigLIP image encoder produces per-frame embeddings (2 fps with keyframe deduplication).",
            "Audio stream: Whisper-large-v3 transcribes with word-level timestamps; bge-m3 embeds 10-second chunks.",
            "Fusion: weighted sum of normalized visual and audio similarities to the question.",
            "Answer: top-K moments + transcript snippets sent to Llama 3.3 (Groq) for synthesis.",
        ],
    )
    add_footer(slide, prs, 5, total)


def slide_bilingual(prs, total):
    slide = blank_slide(prs)
    add_top_bar(slide, prs)
    add_slide_title(slide, prs, "Bilingual Handling")
    add_bullets(
        slide,
        Inches(0.7),
        Inches(1.7),
        prs.slide_width - Inches(1.4),
        Inches(5.0),
        [
            "Whisper-large-v3 auto-detects Arabic and English audio.",
            "bge-m3 retrieves over Arabic and English transcripts natively.",
            "SigLIP is English-only; Arabic questions are translated for visual retrieval and the translation is shown to the user.",
            "Final answer is generated in the language of the question.",
        ],
    )
    add_footer(slide, prs, 6, total)


def slide_xai(prs, total):
    slide = blank_slide(prs)
    add_top_bar(slide, prs)
    add_slide_title(slide, prs, "Explainable AI — 4 Views")
    add_bullets(
        slide,
        Inches(0.7),
        Inches(1.55),
        prs.slide_width - Inches(1.4),
        Inches(4.6),
        [
            "Timeline relevance: per-second visual and audio relevance over the full video, peaks marked.",
            "Top candidate frames: top-K frames as a grid, each labeled with its visual / audio / combined score.",
            "Modality contribution: stacked bar chart per moment — did the answer come from sight or sound?",
            "Grad-CAM: spatial attention heatmap on the top frame, showing which regions matched the question.",
        ],
        font_size=Pt(20),
    )
    add_textbox(
        slide,
        Inches(0.7),
        Inches(6.35),
        prs.slide_width - Inches(1.4),
        Inches(0.5),
        'Together these answer "what, where, when, and why."',
        font_size=Pt(18),
        bold=True,
        color=ACCENT,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, prs, 7, total)


def slide_evaluation(prs, total):
    slide = blank_slide(prs)
    add_top_bar(slide, prs)
    add_slide_title(slide, prs, "Evaluation Methodology")
    add_bullets(
        slide,
        Inches(0.7),
        Inches(1.7),
        prs.slide_width - Inches(1.4),
        Inches(5.0),
        [
            "Annotated triples: (video, question, ground-truth timestamp).",
            "Metrics: Top-1 IoU, Hit@1 / Hit@3 / Hit@5, Mean Temporal Error (seconds).",
            "Note: small annotated set; designed as a methodology, not a leaderboard claim.",
        ],
    )
    add_footer(slide, prs, 8, total)


def slide_tech_stack(prs, total):
    slide = blank_slide(prs)
    add_top_bar(slide, prs)
    add_slide_title(slide, prs, "Tech Stack")

    col_width = (prs.slide_width - Inches(2.0)) / 2
    left_x = Inches(0.7)
    right_x = left_x + col_width + Inches(0.6)

    add_textbox(
        slide,
        left_x,
        Inches(1.6),
        col_width,
        Inches(0.5),
        "Models",
        font_size=Pt(24),
        bold=True,
        color=ACCENT,
    )
    add_bullets(
        slide,
        left_x,
        Inches(2.2),
        col_width,
        Inches(4.5),
        [
            "SigLIP (google/siglip-large-patch16-384)",
            "Whisper-large-v3 (faster-whisper / CTranslate2)",
            "BAAI/bge-m3 dense multilingual retriever",
            "Llama 3.3 70B Versatile (Groq API)",
        ],
        font_size=Pt(18),
    )

    add_textbox(
        slide,
        right_x,
        Inches(1.6),
        col_width,
        Inches(0.5),
        "Tools",
        font_size=Pt(24),
        bold=True,
        color=ACCENT,
    )
    add_bullets(
        slide,
        right_x,
        Inches(2.2),
        col_width,
        Inches(4.5),
        [
            "PyTorch",
            "Hugging Face Transformers",
            "faster-whisper",
            "pytorch-grad-cam",
            "Gradio",
            "Groq API",
            "Google Colab",
        ],
        font_size=Pt(18),
    )
    add_footer(slide, prs, 9, total)


def slide_limitations(prs, total):
    slide = blank_slide(prs)
    add_top_bar(slide, prs)
    add_slide_title(slide, prs, "Limitations & Future Work")
    add_bullets(
        slide,
        Inches(0.7),
        Inches(1.6),
        prs.slide_width - Inches(1.4),
        Inches(5.4),
        [
            "Designed for short videos (< 10 min); longer videos need config tuning.",
            "Whisper diarization not used; multi-speaker accuracy may suffer.",
            "SigLIP is English-trained; Arabic visual queries depend on translation quality.",
            "URL downloads from Colab IPs sometimes fail; direct upload is the reliable path.",
            "Future: end-to-end video-language model comparison; speaker-aware retrieval; finer temporal resolution.",
        ],
        font_size=Pt(20),
    )
    add_footer(slide, prs, 10, total)


def slide_thanks(prs, total):
    slide = blank_slide(prs)
    add_top_bar(slide, prs)
    add_textbox(
        slide,
        Inches(0.6),
        Inches(2.6),
        prs.slide_width - Inches(1.2),
        Inches(1.4),
        "Thank you. Questions?",
        font_size=Pt(54),
        bold=True,
        color=ACCENT,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        Inches(0.6),
        Inches(4.6),
        prs.slide_width - Inches(1.2),
        Inches(0.5),
        "https://github.com/Abdelrahman-Nashaat/Video-Question-Answering-with-Temporal-Grounding",
        font_size=Pt(16),
        color=TEXT_MUTED,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, prs, 11, total)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    total = 11
    slide_title(prs, total)
    slide_problem(prs, total)
    slide_goals(prs, total)
    slide_architecture(prs, total)
    slide_pipeline(prs, total)
    slide_bilingual(prs, total)
    slide_xai(prs, total)
    slide_evaluation(prs, total)
    slide_tech_stack(prs, total)
    slide_limitations(prs, total)
    slide_thanks(prs, total)

    out = Path(__file__).resolve().parent / "presentation.pptx"
    prs.save(out)
    print(f"Wrote {out} ({out.stat().st_size} bytes)")


if __name__ == "__main__":
    build()
